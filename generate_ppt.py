"""
Génère la présentation PowerPoint (.pptx) pour BININGA
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Couleurs ────────────────────────────────────────────────────────────────
ROUGE   = RGBColor(0xC0, 0x00, 0x00)
BLANC   = RGBColor(0xFF, 0xFF, 0xFF)
NOIR    = RGBColor(0x1A, 0x1A, 0x1A)
GRIS    = RGBColor(0x55, 0x55, 0x55)
GRIS_C  = RGBColor(0xF2, 0xF2, 0xF2)
GRIS_F  = RGBColor(0xDD, 0xDD, 0xDD)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # totalement vide

# ── Helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=24, bold=False, italic=False,
             color=NOIR, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return txBox

def add_multiline(slide, lines, x, y, w, h, size=18, bold=False, color=GRIS,
                  align=PP_ALIGN.LEFT, spacing=1.0):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size)
        run.font.bold = line_bold
        run.font.color.rgb = line_color
        run.font.name = 'Calibri'
    return txBox

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}",
             W - Inches(1.2), H - Inches(0.4), Inches(1.0), Inches(0.3),
             size=9, color=GRIS, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 10

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

# Fond blanc cassé
add_rect(slide, 0, 0, W, H, GRIS_C)

# Bande rouge gauche
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

# Bande rouge en haut
add_rect(slide, 0, 0, W, Inches(0.12), ROUGE)

# Bande rouge bas
add_rect(slide, 0, H - Inches(1.4), W, Inches(1.4), ROUGE)

# Nom
add_text(slide, "Ange Aimé Wilfrid",
         Inches(1.0), Inches(1.4), Inches(11.0), Inches(0.9),
         size=28, color=GRIS, align=PP_ALIGN.CENTER)

add_text(slide, "BININGA",
         Inches(1.0), Inches(2.1), Inches(11.0), Inches(1.5),
         size=72, bold=True, color=ROUGE, align=PP_ALIGN.CENTER)

# Séparateur
add_rect(slide, Inches(2.5), Inches(3.55), Inches(8.3), Inches(0.04), GRIS)

# Titre officiel
add_multiline(slide, [
    ("Garde des Sceaux · Ministre de la Justice", False, 16, GRIS),
    ("des Droits Humains et de la Promotion des Peuples Autochtones", False, 14, GRIS),
    ("Député de la 1re Circonscription d'Ewo · PCT", True, 15, NOIR),
], Inches(1.0), Inches(3.65), Inches(11.3), Inches(1.4), align=PP_ALIGN.CENTER)

# Bandeau bas
add_text(slide, "PRÉSENTATION DU SITE WEB OFFICIEL  ·  MARS 2026",
         Inches(0.5), H - Inches(1.1), Inches(12.3), Inches(0.7),
         size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — POURQUOI UN SITE OFFICIEL
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), ROUGE)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Pourquoi un Site Officiel ?",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC, align=PP_ALIGN.LEFT)

slide_number(slide, 2, TOTAL_SLIDES)

points = [
    ("La communication numérique est aujourd'hui indispensable pour tout homme d'État.",
     "Un site officiel crédibilise, informe et rapproche l'élu de ses citoyens."),
    ("Un outil de communication institutionnel moderne",
     "Disponible 24h/24, accessible depuis n'importe quel appareil dans le monde entier."),
    ("Gestion des relations citoyennes",
     "Les citoyens peuvent soumettre des demandes d'audience et des réclamations directement en ligne."),
    ("Visibilité nationale et internationale",
     "Référencement Google optimisé. Présence en ligne digne d'un Garde des Sceaux."),
]

y = Inches(1.55)
for i, (title, desc) in enumerate(points):
    # Numéro
    add_rect(slide, Inches(0.5), y, Inches(0.55), Inches(0.55), ROUGE)
    add_text(slide, str(i+1),
             Inches(0.5), y, Inches(0.55), Inches(0.55),
             size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    # Texte
    add_text(slide, title,
             Inches(1.25), y - Inches(0.02), Inches(11.5), Inches(0.35),
             size=14, bold=True, color=NOIR)
    add_text(slide, desc,
             Inches(1.25), y + Inches(0.3), Inches(11.5), Inches(0.35),
             size=12, color=GRIS)
    y += Inches(1.2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CE QUI A ÉTÉ LIVRÉ (CHIFFRES)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), NOIR)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Le Projet en Chiffres",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC, align=PP_ALIGN.LEFT)

slide_number(slide, 3, TOTAL_SLIDES)

# 4 grandes stats
stats = [
    ("2 725", "lignes de code écrites"),
    ("146 Ko", "de code livré"),
    ("8", "sections sur le site public"),
    ("10", "modules dans le panneau admin"),
]

gap = Inches(0.3)
box_w = (W - Inches(0.5) - gap * 5) / 4
x = Inches(0.5) + gap
y_box = Inches(1.6)
for val, label in stats:
    add_rect(slide, x, y_box, box_w, Inches(1.6), ROUGE)
    add_text(slide, val, x, y_box + Inches(0.15), box_w, Inches(0.85),
             size=40, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y_box + Inches(0.95), box_w, Inches(0.55),
             size=13, color=BLANC, align=PP_ALIGN.CENTER)
    x += box_w + gap

# Tableau fichiers
rows_data = [
    ("Site public  (index.html)", "1 257 lignes", "78 Ko"),
    ("Panneau admin  (admin.html)", "1 096 lignes", "53 Ko"),
    ("Serveur Python  (server.py)", "230 lignes", "10 Ko"),
    ("Base de données  (data.json)", "142 lignes", "5 Ko"),
]
headers = ["Composante", "Lignes", "Poids"]
col_w = [Inches(6.5), Inches(2.5), Inches(2.5)]
x_start = Inches(0.9)
y_tbl = Inches(3.55)
row_h = Inches(0.52)

# Header
x = x_start
for ci, (h, cw) in enumerate(zip(headers, col_w)):
    add_rect(slide, x, y_tbl, cw, row_h, NOIR)
    add_text(slide, h, x + Inches(0.1), y_tbl + Inches(0.08), cw, row_h,
             size=13, bold=True, color=BLANC,
             align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
    x += cw

for ri, row in enumerate(rows_data):
    y_r = y_tbl + row_h * (ri+1)
    x = x_start
    bg = GRIS_C if ri%2==0 else GRIS_F
    for ci, (val, cw) in enumerate(zip(row, col_w)):
        add_rect(slide, x, y_r, cw, row_h, bg)
        c = ROUGE if ci==0 else GRIS
        add_text(slide, val, x + Inches(0.1), y_r + Inches(0.08), cw, row_h,
                 size=12, bold=(ci==0), color=c,
                 align=PP_ALIGN.CENTER if ci>0 else PP_ALIGN.LEFT)
        x += cw

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LE SITE PUBLIC
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), ROUGE)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Le Site Public — 8 Sections",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC)

slide_number(slide, 4, TOTAL_SLIDES)

sections = [
    ("Accueil", "Animation cinématique, slogan, appel à l'action"),
    ("À propos", "Biographie + 3 compteurs animés"),
    ("Engagement", "Valeurs et engagements politiques"),
    ("Parcours", "Ligne du temps interactive"),
    ("Programme", "Priorités pour Ewo et le Congo"),
    ("Galerie", "Slider photos + grille, 100% éditable"),
    ("Actualités", "Article vedette + liste d'actualités"),
    ("Contact & Audiences", "Formulaires avec envoi email réel"),
]

col1 = sections[:4]
col2 = sections[4:]
y = Inches(1.55)
for (sec, desc), (sec2, desc2) in zip(col1, col2):
    add_rect(slide, Inches(0.5), y, Inches(0.06), Inches(0.65), ROUGE)
    add_text(slide, sec, Inches(0.75), y, Inches(5.3), Inches(0.35),
             size=13, bold=True, color=ROUGE)
    add_text(slide, desc, Inches(0.75), y + Inches(0.3), Inches(5.3), Inches(0.35),
             size=11, color=GRIS)

    add_rect(slide, Inches(6.9), y, Inches(0.06), Inches(0.65), ROUGE)
    add_text(slide, sec2, Inches(7.15), y, Inches(5.3), Inches(0.35),
             size=13, bold=True, color=ROUGE)
    add_text(slide, desc2, Inches(7.15), y + Inches(0.3), Inches(5.3), Inches(0.35),
             size=11, color=GRIS)
    y += Inches(1.35)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DESIGN & EXPÉRIENCE UTILISATEUR
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, NOIR)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Design & Expérience Utilisateur",
         Inches(0.5), Inches(0.25), Inches(12.0), Inches(0.85),
         size=30, bold=True, color=BLANC)

slide_number(slide, 5, TOTAL_SLIDES)

ux_items = [
    ("Curseur rouge animé", "Expérience premium — différencie immédiatement du commun des sites"),
    ("Loader animé", "Impression de fluidité et de professionnalisme à l'ouverture"),
    ("Animations au scroll", "Chaque section entre en scène — effet cinématique"),
    ("Responsive 100%", "Parfait sur téléphone, tablette, ordinateur"),
    ("Police Raleway + Lato", "Typographie premium, identité visuelle forte"),
    ("Palette rouge & blanc", "Couleurs institutionnelles, impact politique"),
    ("Menu mobile burger", "Navigation intuitive sur smartphone"),
    ("Smooth scroll", "Navigation fluide entre les sections"),
]

y = Inches(1.4)
for i, (title, desc) in enumerate(ux_items):
    col = 0 if i < 4 else 1
    row = i % 4
    x = Inches(0.6) if col==0 else Inches(7.0)
    y_pos = Inches(1.4) + row * Inches(1.35)
    add_rect(slide, x, y_pos + Inches(0.08), Inches(0.4), Inches(0.4),
             RGBColor(0xC0, 0x00, 0x00))
    add_text(slide, "✓", x, y_pos, Inches(0.4), Inches(0.5),
             size=18, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.55), y_pos, Inches(5.6), Inches(0.35),
             size=13, bold=True, color=BLANC)
    add_text(slide, desc, x + Inches(0.55), y_pos + Inches(0.32), Inches(5.6), Inches(0.45),
             size=11, color=RGBColor(0xAA, 0xAA, 0xAA))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PANNEAU ADMIN (CMS)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), NOIR)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Panneau Admin — CMS Sur Mesure",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC)

slide_number(slide, 6, TOTAL_SLIDES)

# Sous-titre
add_text(slide, "Gérez l'intégralité du site sans aucune compétence technique",
         Inches(0.5), Inches(1.35), Inches(12.0), Inches(0.4),
         size=14, italic=True, color=GRIS)

modules = [
    ("Tableau de bord", "Statistiques en temps réel"),
    ("Éditeur Hero", "Nom, rôle, slogan"),
    ("Éditeur À propos", "Biographie, intro"),
    ("Éditeur Stats", "Les 3 compteurs"),
    ("Éditeur SEO", "Google, référencement"),
    ("Galerie Slider", "Photos, ordre, légendes"),
    ("Galerie Grille", "Ajout / suppression photos"),
    ("Actualités", "Articles, dates, textes"),
    ("Demandes d'audience", "Filtres, statuts, suivi"),
    ("Réclamations", "Gestion citoyenne"),
    ("Messages Contact", "Suivi des messages"),
]

cols = 3
box_w = Inches(3.7)
box_h = Inches(0.85)
gap_x = Inches(0.3)
gap_y = Inches(0.2)
x0 = Inches(0.5)
y0 = Inches(1.95)

for i, (mod, sub) in enumerate(modules):
    col = i % cols
    row = i // cols
    x = x0 + col * (box_w + gap_x)
    y = y0 + row * (box_h + gap_y)
    add_rect(slide, x, y, box_w, box_h, GRIS_C)
    add_rect(slide, x, y, Inches(0.08), box_h, ROUGE)
    add_text(slide, mod, x + Inches(0.2), y + Inches(0.06), box_w - Inches(0.3), Inches(0.4),
             size=12, bold=True, color=ROUGE)
    add_text(slide, sub, x + Inches(0.2), y + Inches(0.42), box_w - Inches(0.3), Inches(0.35),
             size=10, color=GRIS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SÉCURITÉ
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, GRIS_C)
add_rect(slide, 0, 0, W, Inches(1.3), ROUGE)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Sécurité Institutionnelle",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC)

slide_number(slide, 7, TOTAL_SLIDES)

sec_items = [
    ("Authentification serveur", "Le mot de passe est vérifié côté serveur — impossible à contourner depuis le navigateur."),
    ("Token de session dynamique", "Chaque connexion génère un code unique. Accès admin protégé."),
    ("Connexion HTTPS chiffrée", "Toutes les données transitent de façon chiffrée (SSL/TLS natif)."),
    ("Sauvegarde automatique", "À chaque modification, une copie de sauvegarde est créée automatiquement."),
    ("Protection anti-XSS", "Toutes les saisies citoyennes sont neutralisées avant affichage."),
    ("Contrôle des uploads", "Seules les images autorisées. Taille maximale contrôlée côté serveur."),
]

y = Inches(1.55)
for i, (title, desc) in enumerate(sec_items):
    col = 0 if i < 3 else 1
    row = i % 3
    x = Inches(0.5) if col==0 else Inches(7.0)
    y_pos = Inches(1.55) + row * Inches(1.7)
    add_rect(slide, x, y_pos, Inches(5.9), Inches(1.45), BLANC)
    add_rect(slide, x, y_pos, Inches(0.1), Inches(1.45), ROUGE)
    add_text(slide, title, x + Inches(0.25), y_pos + Inches(0.1), Inches(5.5), Inches(0.45),
             size=13, bold=True, color=ROUGE)
    add_text(slide, desc, x + Inches(0.25), y_pos + Inches(0.5), Inches(5.5), Inches(0.75),
             size=11, color=GRIS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — VALEUR DU PROJET
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), NOIR)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Valeur du Projet",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC)

slide_number(slide, 8, TOTAL_SLIDES)

# Équipe nécessaire
add_text(slide, "Dans une agence professionnelle, ce projet aurait nécessité :",
         Inches(0.5), Inches(1.4), Inches(12.0), Inches(0.4),
         size=13, color=GRIS, italic=True)

team = [
    ("Designer UI/UX", "Maquettes & identité visuelle"),
    ("Intégrateur Front-End", "HTML, CSS, animations"),
    ("Développeur JavaScript", "Logique & CMS admin"),
    ("Développeur Back-End", "Serveur, API, sécurité"),
    ("Chef de Projet", "Coordination & livraison"),
]

gap = Inches(0.25)
bw = (W - Inches(0.5) - gap * 6) / 5
x = Inches(0.5) + gap
y_b = Inches(1.95)
for role, sub in team:
    add_rect(slide, x, y_b, bw, Inches(1.6), GRIS_C)
    add_rect(slide, x, y_b, bw, Inches(0.1), ROUGE)
    add_text(slide, role, x + Inches(0.1), y_b + Inches(0.2), bw - Inches(0.2), Inches(0.8),
             size=12, bold=True, color=ROUGE, align=PP_ALIGN.CENTER)
    add_text(slide, sub, x + Inches(0.1), y_b + Inches(0.9), bw - Inches(0.2), Inches(0.55),
             size=10, color=GRIS, align=PP_ALIGN.CENTER)
    x += bw + gap

add_text(slide, "= 4 à 5 personnes · 3 à 6 semaines de travail",
         Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.5),
         size=14, bold=True, color=NOIR, align=PP_ALIGN.CENTER)

# Pricing
add_text(slide, "Estimation comparative de marché :",
         Inches(0.5), Inches(4.35), Inches(12.0), Inches(0.4),
         size=13, bold=True, color=NOIR)

pricing = [
    ("Freelance Congo / Afrique Centrale", "400 000 – 800 000 FCFA"),
    ("Agence web africaine", "800 000 – 2 000 000 FCFA"),
    ("Agence web France / Europe", "4 000 – 10 000 €"),
]

pw = Inches(3.8)
gap_p = Inches(0.25)
x = Inches(0.5) + gap_p
for ctx, price in pricing:
    add_rect(slide, x, Inches(4.85), pw, Inches(1.7), GRIS_C)
    add_rect(slide, x, Inches(4.85), pw, Inches(0.08), ROUGE)
    add_text(slide, ctx, x + Inches(0.15), Inches(4.98), pw - Inches(0.3), Inches(0.6),
             size=11, color=GRIS)
    add_text(slide, price, x + Inches(0.15), Inches(5.55), pw - Inches(0.3), Inches(0.65),
             size=16, bold=True, color=ROUGE, align=PP_ALIGN.CENTER)
    x += pw + gap_p

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PROCHAINES ÉTAPES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, BLANC)
add_rect(slide, 0, 0, W, Inches(1.3), ROUGE)
add_rect(slide, 0, 0, Inches(0.18), H, ROUGE)

add_text(slide, "Prochaines Étapes",
         Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.9),
         size=30, bold=True, color=BLANC)

slide_number(slide, 9, TOTAL_SLIDES)

steps = [
    ("1", "Hébergement", "Mise en ligne sur un serveur officiel avec nom de domaine (ex: bininga.cg)"),
    ("2", "Photos officielles", "Intégration des photos réelles du ministre dans la galerie"),
    ("3", "Formation", "Formation de l'équipe à l'utilisation du panneau admin"),
    ("4", "Maintenance", "Contrat de maintenance mensuel pour mises à jour et sécurité"),
    ("5", "Évolutions", "Ajout de nouvelles fonctionnalités selon les besoins (vidéos, newsletter...)"),
]

y = Inches(1.55)
for num, title, desc in steps:
    add_rect(slide, Inches(0.5), y, Inches(0.7), Inches(0.7), ROUGE)
    add_text(slide, num, Inches(0.5), y, Inches(0.7), Inches(0.7),
             size=22, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.4), y, Inches(11.0), Inches(0.38),
             size=14, bold=True, color=ROUGE)
    add_text(slide, desc, Inches(1.4), y + Inches(0.36), Inches(11.0), Inches(0.38),
             size=12, color=GRIS)
    y += Inches(1.12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, ROUGE)

# Motif subtil
add_rect(slide, W - Inches(4), 0, Inches(4), H, RGBColor(0xA0, 0x00, 0x00))

add_text(slide, "Ange Aimé Wilfrid",
         Inches(0.7), Inches(1.0), Inches(10.0), Inches(0.8),
         size=22, color=BLANC, align=PP_ALIGN.CENTER)

add_text(slide, "BININGA",
         Inches(0.7), Inches(1.7), Inches(10.0), Inches(1.4),
         size=60, bold=True, color=BLANC, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(2.0), Inches(3.1), Inches(9.3), Inches(0.04),
         RGBColor(0xFF, 0xFF, 0xFF))

add_text(slide, '"Pour Ewo. Pour un Congo juste, libre et fort."',
         Inches(0.7), Inches(3.3), Inches(11.9), Inches(0.65),
         size=18, italic=True, color=BLANC, align=PP_ALIGN.CENTER)

add_text(slide, "Un site officiel à la hauteur de l'homme d'État.",
         Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.55),
         size=15, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

add_text(slide, "Mars 2026  ·  République du Congo",
         Inches(0.7), H - Inches(0.9), Inches(11.9), Inches(0.5),
         size=12, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
prs.save("Presentation_BININGA.pptx")
print("✅  Présentation PowerPoint générée : Presentation_BININGA.pptx")
